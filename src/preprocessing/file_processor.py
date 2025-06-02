from pptx import Presentation
import nbformat
import re
import os
import shutil
from pathlib import Path

class FileProcessor:
    @staticmethod
    def process_slide_file(file_path: str, mentor_originals: str, mentor_transcripts: str) -> str:
        file_name = os.path.basename(file_path)
        base_name = os.path.splitext(file_name)[0]
        transcript_path = os.path.join(mentor_transcripts, f"{base_name}.txt")

        # Save original
        original_dest = os.path.join(mentor_originals, file_name)
        if not os.path.exists(original_dest):
            shutil.copy2(file_path, original_dest)

        try:
            prs = Presentation(file_path)
            content = []

            for i, slide in enumerate(prs.slides):
                content.append(f"=== Slide {i+1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        cleaned_text = re.sub(r'\s+', ' ', shape.text.strip())
                        content.append(cleaned_text)
                content.append("")

            with open(transcript_path, 'w', encoding='utf-8') as f:
                f.write("\n".join(content))

            with open(transcript_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error processing presentation: {str(e)}")
            return ""

    @staticmethod
    def process_notebook_file(file_path: str, mentor_originals: str, mentor_transcripts: str) -> str:
        file_name = os.path.basename(file_path)
        base_name = os.path.splitext(file_name)[0]
        transcript_path = os.path.join(mentor_transcripts, f"{base_name}.txt")

        # Save original
        original_dest = os.path.join(mentor_originals, file_name)
        if not os.path.exists(original_dest):
            shutil.copy2(file_path, original_dest)

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                nb = nbformat.read(f, as_version=4)

            content = []
            for cell in nb.cells:
                if cell.cell_type == 'code':
                    content.append("code:")
                    content.append(cell.source.strip())
                    content.append("----")
                elif cell.cell_type == 'markdown':
                    content.append("text:")
                    cleaned_text = cell.source.strip()
                    cleaned_text = re.sub(r'#+\s*', '', cleaned_text)
                    cleaned_text = re.sub(r'\*{1,2}(.*?)\*{1,2}', r'\1', cleaned_text)
                    cleaned_text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', cleaned_text)
                    content.append(cleaned_text)
                    content.append("----")

            with open(transcript_path, 'w', encoding='utf-8') as f:
                f.write("\n".join(content))

            with open(transcript_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error processing notebook: {str(e)}")
            return ""